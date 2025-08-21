import io
from pptx import Presentation
from pptx.util import Inches, Pt

def build_presentation_from_plan(plan, template_stream):
    """Build a pptx bytes object from the plan using the uploaded template.
    plan: list of slide definitions (title, type, bullets, notes)
    template_stream: BytesIO of uploaded template
    Returns bytes of the generated PPTX.
    """
    # Load template (if invalid, start fresh)
    try:
        prs = Presentation(template_stream)
    except Exception:
        prs = Presentation()

    # Simple mapping: use layout 1 (title + content) if available, else blank slide
    layout_title_and_content = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]

    # We'll append new slides to the end of the presentation
    for s in plan:
        slide_type = s.get('type', 'title_and_bullets')
        if slide_type in ('title',):
            layout = prs.slide_layouts[0]
        else:
            layout = layout_title_and_content

        slide = prs.slides.add_slide(layout)
        # Title
        try:
            slide.shapes.title.text = s.get('title','')
        except Exception:
            # If no title placeholder, create a textbox
            from pptx.util import Pt
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
            tf = txBox.text_frame
            tf.text = s.get('title','')

        # Content / bullets
        bullets = s.get('bullets', [])
        # Find first placeholder that has a text_frame and is not the title
        body_shape = None
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and not shape == slide.shapes.title:
                body_shape = shape
                break
        if body_shape is None:
            # create a textbox
            body_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(4))

        tf = body_shape.text_frame
        # clear existing paragraphs
        tf.clear()
        for i, b in enumerate(bullets):
            p = tf.add_paragraph()
            p.text = b
            p.level = 0

        # Notes
        notes = s.get('notes','')
        try:
            slide.notes_slide.notes_text_frame.text = notes
        except Exception:
            pass

    # Save to bytes
    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out.getvalue()
